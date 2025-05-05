import os
import json
from pathlib import Path
from typing import Dict, Tuple, List, Optional
import datetime

import pandas as pd
import numpy as np
import streamlit as st
import joblib
from sklearn.ensemble import RandomForestRegressor
import matplotlib.pyplot as plt

try:
    from fpdf import FPDF  # generación rápida de PDF
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

# PowerPoint export
try:
    from pptx import Presentation
    from pptx.util import Inches
    PPT_AVAILABLE = True
except ImportError:
    PPT_AVAILABLE = False

# Estilos CSS para mejorar la interfaz
CUSTOM_CSS = """
<style>
    .stApp {
        max-width: 1200px;
        margin: 0 auto;
    }
    .main .block-container {
        padding-top: 2rem;
    }
    h1, h2, h3 {
        color: #1E3A8A;
    }
    .stMetric {
        background-color: #f0f5ff;
        padding: 10px;
        border-radius: 5px;
        border-left: 3px solid #1E3A8A;
    }
    .stAlert {
        border-radius: 4px;
    }
    /* Mejora de barras de progreso en textos */
    pre {
        font-size: 1.2em !important;
    }
    .css-1kyxreq {
        justify-content: center;
    }
</style>
"""

# ============================================================
# Configuración global de Streamlit
# ============================================================

st.set_page_config(
    page_title="Motor de Decisión Educativo",
    page_icon="📘",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ============================================================
# Utilidades de gestión de marcas y datos
# ============================================================

DATA_DIR = Path("data")
DATA_DIR.mkdir(exist_ok=True)

# Definir las marcas predefinidas
MARCAS_PREDEFINIDAS = ["GR", "PG", "ADV", "GMBA", "WZ", "AJA"]

# Configuración de estructura de directorios
REPORTE_DIR_NAME = "reportes"
HISTORICO_DIR_NAME = "historico"
ACTUAL_DIR_NAME = "actual"
LEADS_DIR_NAME = "leads"
MATRICULAS_DIR_NAME = "matriculas"
PLAN_FILE = "plan_actual.csv"
HIST_FILE = "historico.csv"
MODEL_FILE = "modelo_rf.joblib"

# Versión del sistema
VERSION = "2.0.0"

# Tooltip ayudas contextuales
TOOLTIPS = {
    "cpa": "CPA (Costo Por Adquisición) es el costo promedio de obtener una matrícula.",
    "cpl": "CPL (Costo Por Lead) es el costo promedio de obtener un lead potencial.",
    "leads": "Prospecto interesado en un programa educativo.",
    "prediccion_ml": "Predicción basada en Machine Learning (RandomForest) con intervalos de confianza.",
    "anomalia": "Valor que se desvía significativamente del resto (> 3 desviaciones estándar).",
    "atribucion": "Modelo que determina cómo asignar el mérito a cada canal en la conversión de un lead."
}

def get_brand_path(brand: str) -> Path:
    """Devuelve la ruta del directorio de la marca y crea la estructura completa si no existe."""
    # Directorio principal de la marca
    brand_path = DATA_DIR / brand
    brand_path.mkdir(exist_ok=True)
    
    # Subdirectorios principales
    (brand_path / REPORTE_DIR_NAME).mkdir(exist_ok=True)
    (brand_path / HISTORICO_DIR_NAME).mkdir(exist_ok=True)
    (brand_path / ACTUAL_DIR_NAME).mkdir(exist_ok=True)
    
    # Subdirectorios de histórico
    (brand_path / HISTORICO_DIR_NAME / LEADS_DIR_NAME).mkdir(exist_ok=True)
    (brand_path / HISTORICO_DIR_NAME / MATRICULAS_DIR_NAME).mkdir(exist_ok=True)
    
    # Subdirectorios para año actual
    current_year = str(datetime.datetime.now().year)
    (brand_path / HISTORICO_DIR_NAME / LEADS_DIR_NAME / current_year).mkdir(exist_ok=True)
    (brand_path / HISTORICO_DIR_NAME / MATRICULAS_DIR_NAME / current_year).mkdir(exist_ok=True)
    
    return brand_path


def save_dataframe(df: pd.DataFrame, path: Path):
    """Guarda un DataFrame en CSV."""
    df.to_csv(path, index=False)


def load_dataframe(path: Path) -> pd.DataFrame:
    """Carga un CSV a DataFrame, si existe."""
    if path.exists():
        return pd.read_csv(path)
    return pd.DataFrame()


def get_weekly_filename(brand: str, week: int, year: int, data_type: str) -> str:
    """Genera un nombre de archivo estándar para datos semanales."""
    return f"{brand}_{data_type}_S{week:02d}_{year}.csv"


def save_weekly_data(df: pd.DataFrame, brand: str, week: int, year: int, data_type: str):
    """Guarda datos semanales en el directorio apropiado."""
    brand_path = get_brand_path(brand)
    
    if data_type == "leads":
        target_dir = brand_path / HISTORICO_DIR_NAME / LEADS_DIR_NAME / str(year)
    elif data_type == "matriculas":
        target_dir = brand_path / HISTORICO_DIR_NAME / MATRICULAS_DIR_NAME / str(year)
    else:
        target_dir = brand_path / HISTORICO_DIR_NAME
    
    target_dir.mkdir(exist_ok=True)
    filename = get_weekly_filename(brand, week, year, data_type)
    save_dataframe(df, target_dir / filename)


def load_all_historical_data(brand: str, data_type: str) -> pd.DataFrame:
    """Carga todos los datos históricos de un tipo específico."""
    brand_path = get_brand_path(brand)
    
    if data_type == "leads":
        base_dir = brand_path / HISTORICO_DIR_NAME / LEADS_DIR_NAME
    elif data_type == "matriculas":
        base_dir = brand_path / HISTORICO_DIR_NAME / MATRICULAS_DIR_NAME
    else:
        # Para otros tipos, buscar en el directorio principal
        if (brand_path / f"{data_type}_historico.csv").exists():
            return load_dataframe(brand_path / f"{data_type}_historico.csv")
        return pd.DataFrame()
    
    # Buscar en todos los subdirectorios de años
    dfs = []
    for year_dir in base_dir.iterdir():
        if year_dir.is_dir():
            for file in year_dir.glob(f"{brand}_{data_type}_*.csv"):
                try:
                    df = load_dataframe(file)
                    dfs.append(df)
                except Exception as e:
                    st.warning(f"Error al cargar {file}: {str(e)}")
    
    # Combinar todos los DataFrames
    if dfs:
        return pd.concat(dfs, ignore_index=True)
    return pd.DataFrame()


def check_for_duplicates(df: pd.DataFrame, id_column: str = "ID") -> pd.DataFrame:
    """Verifica si hay IDs duplicados en el DataFrame."""
    if id_column in df.columns:
        duplicates = df[df.duplicated(subset=[id_column], keep=False)]
        return duplicates
    return pd.DataFrame()


def validate_data_structure(df: pd.DataFrame, data_type: str) -> dict:
    """Valida la estructura del DataFrame según el tipo de datos."""
    results = {"valid": True, "missing_columns": [], "message": ""}
    
    required_columns = {
        "planificacion": ["fecha", "marca", "canal", "presupuesto", "leads_estimados", "objetivo_matriculas"],
        "historico": ["fecha", "marca", "canal", "leads", "matriculas", "inversion"],
        "leads": ["ID", "fecha_generacion", "canal", "programa", "marca", "estado"],
        "matriculas": ["ID", "fecha_matricula", "canal", "marca", "programa"]
    }
    
    if data_type in required_columns:
        for col in required_columns[data_type]:
            if col not in df.columns:
                results["missing_columns"].append(col)
                results["valid"] = False
        
        if not results["valid"]:
            results["message"] = f"Faltan columnas requeridas: {', '.join(results['missing_columns'])}"
    else:
        results["valid"] = False
        results["message"] = f"Tipo de datos '{data_type}' no reconocido."
    
    return results


# ============================================================
# Funciones de Machine Learning
# ============================================================

def train_or_load_model(df_hist: pd.DataFrame, brand_path: Path):
    """Entrena un modelo RandomForest o carga uno existente."""
    model_path = brand_path / MODEL_FILE
    
    # Si no hay suficientes datos, no se puede entrenar un modelo
    if len(df_hist) < 10 or "leads" not in df_hist.columns or "matriculas" not in df_hist.columns:
        return None
    
    # Si ya existe un modelo, cargarlo
    if model_path.exists():
        try:
            return joblib.load(model_path)
        except Exception as e:
            st.warning(f"Error al cargar modelo existente: {str(e)}")
    
    # Entrenar un nuevo modelo
    try:
        X = df_hist[["leads", "inversion"]]
        y = df_hist["matriculas"]
        
        model = RandomForestRegressor(n_estimators=100, random_state=42)
        model.fit(X, y)
        
        # Guardar el modelo
        joblib.dump(model, model_path)
        return model
    except Exception as e:
        st.warning(f"Error al entrenar modelo: {str(e)}")
        return None


def predict_matriculas_interval(model, df_future: pd.DataFrame, 
                              confidence_level: float = 0.95) -> Tuple[np.ndarray, Tuple[np.ndarray, np.ndarray]]:
    """
    Realiza predicciones y genera intervalos de confianza.
    
    Args:
        model: Modelo RandomForest entrenado
        df_future: DataFrame con columnas 'leads' e 'inversion'
        confidence_level: Nivel de confianza para el intervalo (0-1)
        
    Returns:
        Tupla con (predicciones, (límite_inferior, límite_superior))
    """
    if model is None:
        return np.zeros(len(df_future)), (np.zeros(len(df_future)), np.zeros(len(df_future)))
    
    X_future = df_future[["leads", "inversion"]]
    
    # Predicción central
    preds = model.predict(X_future)
    
    # Calcular intervalos usando los árboles individuales del RandomForest
    tree_preds = np.array([tree.predict(X_future) for tree in model.estimators_])
    
    # Desviación estándar basada en predicciones de árboles individuales
    std_dev = np.std(tree_preds, axis=0)
    
    # Factor z para el nivel de confianza (aproximación)
    z = 1.96  # para 95% de confianza
    
    lower_bound = preds - z * std_dev
    upper_bound = preds + z * std_dev
    
    # Asegurar que los límites no sean negativos
    lower_bound = np.maximum(lower_bound, 0)
    
    return preds, (lower_bound, upper_bound)


# ============================================================
# Módulo de Atribución Multicanal
# ============================================================

def calcular_atribucion(df_leads: pd.DataFrame, df_matriculas: pd.DataFrame, modelo: str = "ultimo_clic") -> pd.DataFrame:
    """
    Calcula la atribución de matrículas a canales según diferentes modelos.
    
    Args:
        df_leads: DataFrame con leads individuales (debe tener ID, canal)
        df_matriculas: DataFrame con matrículas individuales (debe tener ID)
        modelo: Modelo de atribución a utilizar
            - "ultimo_clic": Atribuye al último canal de contacto
            - "primer_clic": Atribuye al primer canal de contacto
            - "lineal": Distribuye equitativamente entre todos los canales de contacto
            - "tiempo": Atribución con decaimiento temporal
            - "posicional": Mayor peso a primer y último contacto
            - "shapley": Usa cálculo de valor Shapley (computacionalmente intensivo)
    
    Returns:
        DataFrame con la atribución por canal
    """
    # Validar inputs
    if df_leads.empty or df_matriculas.empty:
        return pd.DataFrame()
    
    if "ID" not in df_leads.columns or "ID" not in df_matriculas.columns:
        st.warning("Se requiere la columna 'ID' en ambos DataFrames para calcular atribución")
        return pd.DataFrame()
    
    if "canal" not in df_leads.columns:
        st.warning("Se requiere la columna 'canal' en el DataFrame de leads")
        return pd.DataFrame()
    
    # Obtener IDs de leads que se convirtieron en matrículas
    leads_convertidos = df_leads[df_leads["ID"].isin(df_matriculas["ID"])]
    
    # Si no hay secuencia de canales, usar canal único
    if "secuencia_canales" not in leads_convertidos.columns:
        # Crear secuencia basada en canal único
        leads_convertidos["secuencia_canales"] = leads_convertidos["canal"]
    
    # Implementar diferentes modelos de atribución
    resultados = {}
    
    if modelo == "ultimo_clic":
        # Para cada matrícula, atribuir 100% al último canal
        for id_lead, grupo in leads_convertidos.groupby("ID"):
            ultimo_canal = grupo.iloc[-1]["canal"]
            if ultimo_canal in resultados:
                resultados[ultimo_canal] += 1
            else:
                resultados[ultimo_canal] = 1
    
    elif modelo == "primer_clic":
        # Para cada matrícula, atribuir 100% al primer canal
        for id_lead, grupo in leads_convertidos.groupby("ID"):
            primer_canal = grupo.iloc[0]["canal"]
            if primer_canal in resultados:
                resultados[primer_canal] += 1
            else:
                resultados[primer_canal] = 1
    
    elif modelo == "lineal":
        # Distribuir equitativamente entre todos los canales
        for id_lead, grupo in leads_convertidos.groupby("ID"):
            canales = grupo["canal"].unique()
            valor_por_canal = 1.0 / len(canales)
            for canal in canales:
                if canal in resultados:
                    resultados[canal] += valor_por_canal
                else:
                    resultados[canal] = valor_por_canal
    
    elif modelo == "tiempo":
        # Atribución con decaimiento temporal (más recientes tienen más peso)
        for id_lead, grupo in leads_convertidos.groupby("ID"):
            # Ordenar por fecha si está disponible
            if "fecha_generacion" in grupo.columns:
                grupo = grupo.sort_values("fecha_generacion")
            
            # Aplicar decaimiento exponencial
            total_contactos = len(grupo)
            for i, (_, contacto) in enumerate(grupo.iterrows()):
                # Peso exponencial: los últimos tienen más peso
                peso = 2 ** (i / (total_contactos - 1)) if total_contactos > 1 else 1
                canal = contacto["canal"]
                
                if canal in resultados:
                    resultados[canal] += peso
                else:
                    resultados[canal] = peso
        
        # Normalizar para que la suma total sea igual al número de matrículas
        total_peso = sum(resultados.values())
        total_matriculas = len(leads_convertidos["ID"].unique())
        for canal in resultados:
            resultados[canal] = (resultados[canal] / total_peso) * total_matriculas
    
    elif modelo == "posicional":
        # 40% primer contacto, 40% último contacto, 20% distribuido entre el resto
        for id_lead, grupo in leads_convertidos.groupby("ID"):
            # Ordenar por fecha si está disponible
            if "fecha_generacion" in grupo.columns:
                grupo = grupo.sort_values("fecha_generacion")
            
            canales = grupo["canal"].tolist()
            if len(canales) == 1:
                # Si solo hay un canal, recibe el 100%
                canal = canales[0]
                if canal in resultados:
                    resultados[canal] += 1
                else:
                    resultados[canal] = 1
            else:
                # Primer contacto (40%)
                primer_canal = canales[0]
                if primer_canal in resultados:
                    resultados[primer_canal] += 0.4
                else:
                    resultados[primer_canal] = 0.4
                
                # Último contacto (40%)
                ultimo_canal = canales[-1]
                if ultimo_canal in resultados:
                    resultados[ultimo_canal] += 0.4
                else:
                    resultados[ultimo_canal] = 0.4
                
                # Contactos intermedios (20% distribuido)
                canales_intermedios = canales[1:-1]
                if canales_intermedios:
                    valor_por_canal = 0.2 / len(canales_intermedios)
                    for canal in canales_intermedios:
                        if canal in resultados:
                            resultados[canal] += valor_por_canal
                        else:
                            resultados[canal] = valor_por_canal
                else:
                    # Si no hay canales intermedios, dividir el 20% entre primer y último
                    resultados[primer_canal] += 0.1
                    resultados[ultimo_canal] += 0.1
    
    elif modelo == "shapley":
        # Implementación simplificada del valor Shapley
        # Esta es una aproximación ya que el verdadero cálculo de Shapley es computacionalmente intensivo
        canales_unicos = df_leads["canal"].unique()
        total_matriculas = len(leads_convertidos["ID"].unique())
        
        # Inicializar resultados
        for canal in canales_unicos:
            resultados[canal] = 0
        
        # Para cada canal, calcular su contribución marginal
        for canal in canales_unicos:
            # Leads que pasaron por este canal
            leads_con_canal = df_leads[df_leads["canal"] == canal]["ID"].unique()
            
            # Matrículas que provienen de leads que pasaron por este canal
            matriculas_con_canal = df_matriculas[df_matriculas["ID"].isin(leads_con_canal)]
            
            # Matrículas sin este canal
            matriculas_sin_canal = df_matriculas[~df_matriculas["ID"].isin(leads_con_canal)]
            
            # Calcular tasas de conversión con y sin este canal
            tasa_con_canal = len(matriculas_con_canal) / len(leads_con_canal) if len(leads_con_canal) > 0 else 0
            leads_sin_canal = df_leads[~df_leads["ID"].isin(leads_con_canal)]["ID"].unique()
            tasa_sin_canal = len(matriculas_sin_canal) / len(leads_sin_canal) if len(leads_sin_canal) > 0 else 0
            
            # La contribución es proporcional a la diferencia de tasas
            contribucion = max(0, tasa_con_canal - tasa_sin_canal)
            
            # Normalizar
            resultados[canal] = contribucion
        
        # Normalizar para que la suma sea igual al total de matrículas
        if sum(resultados.values()) > 0:
            factor = total_matriculas / sum(resultados.values())
            for canal in resultados:
                resultados[canal] *= factor
    
    else:
        st.warning(f"Modelo de atribución '{modelo}' no implementado")
        return pd.DataFrame()
    
    # Convertir resultados a DataFrame
    df_atribucion = pd.DataFrame({
        "canal": list(resultados.keys()),
        "atribucion": list(resultados.values())
    })
    
    # Calcular porcentaje de atribución
    total_atribucion = df_atribucion["atribucion"].sum()
    df_atribucion["porcentaje"] = df_atribucion["atribucion"] / total_atribucion * 100
    
    # Ordenar por atribución descendente
    df_atribucion = df_atribucion.sort_values("atribucion", ascending=False).reset_index(drop=True)
    
    return df_atribucion


def comparar_modelos_atribucion(df_leads: pd.DataFrame, df_matriculas: pd.DataFrame) -> Dict[str, pd.DataFrame]:
    """
    Compara diferentes modelos de atribución y retorna los resultados.
    
    Args:
        df_leads: DataFrame con leads individuales
        df_matriculas: DataFrame con matrículas individuales
    
    Returns:
        Diccionario con los resultados de cada modelo de atribución
    """
    modelos = ["ultimo_clic", "primer_clic", "lineal", "tiempo", "posicional", "shapley"]
    resultados = {}
    
    for modelo in modelos:
        df_atribucion = calcular_atribucion(df_leads, df_matriculas, modelo)
        resultados[modelo] = df_atribucion
    
    return resultados


# ============================================================
# UI auxiliar
# ============================================================


def sidebar_brand_selector() -> str:
    """Permite seleccionar o crear una marca desde la barra lateral."""
    st.sidebar.header("🎯 Selección de marca")
    
    # Mostrar todas las marcas existentes
    existing_brands = [d.name for d in DATA_DIR.iterdir() if d.is_dir()]
    
    # Inicializar la marca actual en la sesión
    if "current_brand" not in st.session_state:
        st.session_state.current_brand = existing_brands[0] if existing_brands else ""

    # Opciones de selección: Nueva marca, Marcas predefinidas, o Marcas existentes
    brand_option = st.sidebar.selectbox(
        "Marca", 
        options=["<Nueva marca>", "<Marcas predefinidas>"] + existing_brands, 
        index=0 if not existing_brands else existing_brands.index(st.session_state.current_brand) + 2 if st.session_state.current_brand in existing_brands else 0
    )

    if brand_option == "<Nueva marca>":
        new_brand = st.sidebar.text_input("Nombre de la nueva marca")
        if new_brand:
            st.session_state.current_brand = new_brand.strip()
    elif brand_option == "<Marcas predefinidas>":
        # Permitir seleccionar marcas predefinidas
        selected_brand = st.sidebar.selectbox(
            "Seleccionar marca predefinida:", 
            MARCAS_PREDEFINIDAS,
            index=0
        )
        if selected_brand:
            st.session_state.current_brand = selected_brand
            # Crear estructura de directorios para la marca predefinida
            get_brand_path(selected_brand)
    else:
        st.session_state.current_brand = brand_option
    
    # Mostrar la configuración avanzada
    if st.sidebar.checkbox("Configuración avanzada"):
        st.sidebar.write("Opciones de carga de datos:")
        if "carga_acumulativa" not in st.session_state:
            st.session_state.carga_acumulativa = True
        
        st.session_state.carga_acumulativa = st.sidebar.radio(
            "Modo de carga:",
            ["Acumulativa (actualizar archivo principal)", "Semanal (archivos separados)"],
            index=0 if st.session_state.carga_acumulativa else 1
        ) == "Acumulativa (actualizar archivo principal)"
    
    return st.session_state.current_brand


# ============================================================
# Lógica de carga de datos
# ============================================================

def load_data_ui(brand: str) -> Tuple[pd.DataFrame, pd.DataFrame]:
    """Interfaz para cargar datos históricos y de la convocatoria actual."""
    st.subheader("📥 Carga de datos")
    
    # Obtener la ruta a los directorios de la marca
    brand_path = get_brand_path(brand)
    
    # Configurar camino de archivos
    plan_path = brand_path / PLAN_FILE
    hist_path = brand_path / HIST_FILE
    
    # Directorio actual para datos agregados
    actual_dir = brand_path / ACTUAL_DIR_NAME
    leads_path = actual_dir / "leads_actual.csv"
    matriculas_path = actual_dir / "matriculas_actual.csv"
    inversion_path = actual_dir / "inversion_actual.csv"
    
    # Configuración de la convocatoria
    st.subheader("⚙️ Configuración de convocatoria")
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        semana_actual = st.number_input("Semana actual del año:", min_value=1, max_value=52, value=datetime.datetime.now().isocalendar()[1])
    with col2:
        year_actual = st.number_input("Año:", min_value=2020, max_value=2030, value=datetime.datetime.now().year)
    with col3:
        duracion_convocatoria = st.number_input("Duración total (semanas):", min_value=1, max_value=26, value=13)
    with col4:
        semanas_restantes = st.number_input("Semanas restantes:", min_value=0, max_value=26, value=max(0, duracion_convocatoria - (datetime.datetime.now().isocalendar()[1] % duracion_convocatoria)))
    
    # Guardar configuración en session_state
    if "config_convocatoria" not in st.session_state:
        st.session_state.config_convocatoria = {}
    
    st.session_state.config_convocatoria = {
        "semana_actual": semana_actual,
        "year_actual": year_actual,
        "duracion_convocatoria": duracion_convocatoria,
        "semanas_restantes": semanas_restantes,
        "progreso": (duracion_convocatoria - semanas_restantes) / duracion_convocatoria
    }
    
    # Mostrar modo de carga actual
    st.info(f"Modo de carga: {'Acumulativa' if st.session_state.get('carga_acumulativa', True) else 'Semanal (archivos separados)'}")

    # Sección de carga de archivos - Pestañas para diferentes tipos de datos
    tabs = st.tabs(["Planificación", "Datos Agregados", "Leads Individuales", "Matrículas Individuales"])
    
    # Tab 1: Planificación
    with tabs[0]:
        st.markdown("### Planificación de campaña")
        
        if plan_path.exists():
            df_plan = load_dataframe(plan_path)
            st.success("Planificación cargada")
            st.dataframe(df_plan.head())
        else:
            df_plan = pd.DataFrame()
            st.info("No hay planificación cargada")

        plan_file = st.file_uploader("Subir planificación (CSV/Excel)", type=["csv", "xlsx"], key="plan")
        
        # Añadir ejemplo descargable para planificación
        ejemplo_plan = Path("datos/plantillas/ejemplo_planificacion.csv")
        if ejemplo_plan.exists():
            with open(ejemplo_plan, "r") as f:
                st.download_button(
                    label="Descargar ejemplo de planificación",
                    data=f,
                    file_name="ejemplo_planificacion.csv",
                    mime="text/csv",
                    help="Descarga un archivo CSV de ejemplo para subir como planificación",
                    key="plan_ejemplo"
                )
        
        if plan_file is not None:
            try:
                if plan_file.name.endswith('.csv'):
                    df_plan = pd.read_csv(plan_file)
                else:
                    df_plan = pd.read_excel(plan_file)
                
                # Validar estructura
                validation = validate_data_structure(df_plan, "planificacion")
                if validation["valid"]:
                    # Asegurarse de que la marca en el archivo coincida
                    if "marca" in df_plan.columns:
                        df_plan["marca"] = brand
                    save_dataframe(df_plan, plan_path)
                    st.success("Planificación guardada correctamente")
                else:
                    st.error(validation["message"])
                    st.info("Columnas requeridas: fecha, marca, canal, presupuesto, leads_estimados, objetivo_matriculas")
            except Exception as e:
                st.error(f"Error al cargar planificación: {str(e)}")
    
    # Tab 2: Datos Agregados
    with tabs[1]:
        st.markdown("### Datos Agregados Semanales")
        
        # Crear formulario para carga de datos agregados
        with st.form("form_datos_agregados"):
            st.write("Ingrese los datos agregados de la semana")
            
            col1, col2 = st.columns(2)
            with col1:
                fecha = st.date_input("Fecha", value=datetime.datetime.now())
                canal = st.selectbox("Canal", options=["Facebook", "Instagram", "Google", "Email", "Orgánico", "Otro"])
                leads = st.number_input("Leads generados", min_value=0, value=0)
            
            with col2:
                matriculas = st.number_input("Matrículas confirmadas", min_value=0, value=0)
                inversion = st.number_input("Inversión ($)", min_value=0.0, value=0.0)
                observaciones = st.text_input("Observaciones")
            
            submitted = st.form_submit_button("Guardar datos de la semana")
            
            if submitted:
                # Crear DataFrame con los datos agregados
                nuevo_registro = pd.DataFrame({
                    "fecha": [fecha.strftime("%Y-%m-%d")],
                    "marca": [brand],
                    "canal": [canal],
                    "leads": [leads],
                    "matriculas": [matriculas],
                    "inversion": [inversion],
                    "observaciones": [observaciones]
                })
                
                # Modo acumulativo: Actualizar archivo principal
                if st.session_state.get("carga_acumulativa", True):
                    # Cargar histórico existente si existe
                    df_hist = load_dataframe(hist_path) if hist_path.exists() else pd.DataFrame()
                    
                    # Combinar con nuevos datos
                    if not df_hist.empty:
                        df_hist = pd.concat([df_hist, nuevo_registro], ignore_index=True)
                    else:
                        df_hist = nuevo_registro
                    
                    # Guardar histórico actualizado
                    save_dataframe(df_hist, hist_path)
                    st.success(f"Datos agregados al histórico (total: {len(df_hist)} registros)")
                
                # Modo semanal: Guardar archivo separado
                else:
                    try:
                        # Guardar archivo semanal
                        save_weekly_data(nuevo_registro, brand, semana_actual, year_actual, "historico")
                        st.success(f"Datos guardados en archivo semanal (Semana {semana_actual}, {year_actual})")
                    except Exception as e:
                        st.error(f"Error al guardar archivo semanal: {str(e)}")
        
        # Permitir también carga por archivo
        st.markdown("### O cargue un archivo")
        hist_file = st.file_uploader("Subir histórico agregado (CSV/Excel)", type=["csv", "xlsx"], key="historico")
        
        if hist_file is not None:
            try:
                if hist_file.name.endswith('.csv'):
                    df_hist_new = pd.read_csv(hist_file)
                else:
                    df_hist_new = pd.read_excel(hist_file)
                
                # Validar estructura
                validation = validate_data_structure(df_hist_new, "historico")
                if validation["valid"]:
                    # Asegurarse que la marca es correcta
                    if "marca" in df_hist_new.columns:
                        df_hist_new["marca"] = brand
                    
                    # Guardar según el modo seleccionado
                    if st.session_state.get("carga_acumulativa", True):
                        # Cargar histórico existente si existe
                        df_hist = load_dataframe(hist_path) if hist_path.exists() else pd.DataFrame()
                        
                        # Combinar con nuevos datos
                        if not df_hist.empty:
                            df_hist = pd.concat([df_hist, df_hist_new], ignore_index=True)
                        else:
                            df_hist = df_hist_new
                        
                        # Guardar histórico actualizado
                        save_dataframe(df_hist, hist_path)
                        st.success(f"Datos agregados al histórico (total: {len(df_hist)} registros)")
                    else:
                        # Guardar como archivo semanal
                        save_weekly_data(df_hist_new, brand, semana_actual, year_actual, "historico")
                        st.success(f"Datos guardados en archivo semanal (Semana {semana_actual}, {year_actual})")
                else:
                    st.error(validation["message"])
                    st.info("Columnas requeridas: fecha, marca, canal, leads, matriculas, inversion")
            except Exception as e:
                st.error(f"Error al cargar histórico: {str(e)}")
        
        # Mostrar histórico actual
        if hist_path.exists():
            with st.expander("Ver histórico actual"):
                df_hist = load_dataframe(hist_path)
                st.dataframe(df_hist)
    
    # Tab 3: Leads Individuales
    with tabs[2]:
        st.markdown("### Datos Individuales de Leads")
        
        # Cargar leads actuales si existen
        if leads_path.exists():
            df_leads_act = load_dataframe(leads_path)
            st.success(f"Leads actuales cargados: {len(df_leads_act)} registros")
            with st.expander("Ver primeros registros"):
                st.dataframe(df_leads_act.head(5))
        else:
            df_leads_act = pd.DataFrame()
            st.info("No hay leads actuales cargados")
        
        # Subir archivo de leads
        leads_file = st.file_uploader("Subir archivo de leads (CSV/Excel)", type=["csv", "xlsx"], key="leads")
        
        if leads_file is not None:
            try:
                if leads_file.name.endswith('.csv'):
                    df_leads_new = pd.read_csv(leads_file)
                else:
                    df_leads_new = pd.read_excel(leads_file)
                
                # Validar estructura
                validation = validate_data_structure(df_leads_new, "leads")
                if validation["valid"]:
                    # Verificar duplicados
                    duplicados = check_for_duplicates(df_leads_new)
                    if not duplicados.empty:
                        st.warning(f"Se encontraron {len(duplicados)} IDs duplicados en el archivo.")
                        with st.expander("Ver duplicados"):
                            st.dataframe(duplicados)
                    
                    # Asegurarse que la marca es correcta
                    if "marca" in df_leads_new.columns:
                        df_leads_new["marca"] = brand
                    
                    # Guardar según el modo seleccionado
                    if st.session_state.get("carga_acumulativa", True):
                        # Combinar con leads actuales si existen
                        if not df_leads_act.empty:
                            # Eliminar duplicados basados en ID
                            if "ID" in df_leads_act.columns and "ID" in df_leads_new.columns:
                                df_leads_act = df_leads_act[~df_leads_act["ID"].isin(df_leads_new["ID"])]
                            df_leads_act = pd.concat([df_leads_act, df_leads_new], ignore_index=True)
                        else:
                            df_leads_act = df_leads_new
                        
                        # Guardar leads actualizados
                        save_dataframe(df_leads_act, leads_path)
                        st.success(f"Leads actualizados (total: {len(df_leads_act)} registros)")
                    else:
                        # Guardar como archivo semanal
                        save_weekly_data(df_leads_new, brand, semana_actual, year_actual, "leads")
                        st.success(f"Leads guardados en archivo semanal (Semana {semana_actual}, {year_actual})")
                else:
                    st.error(validation["message"])
                    st.info("Columnas requeridas: ID, fecha_generacion, canal, programa, marca, estado")
            except Exception as e:
                st.error(f"Error al cargar leads: {str(e)}")
        
        # Ejemplo de leads
        ejemplo_leads = Path("datos/plantillas/ejemplo_leads.csv")
        if ejemplo_leads.exists():
            with open(ejemplo_leads, "r") as f:
                st.download_button(
                    label="Ejemplo de estructura de leads",
                    data=f,
                    file_name="ejemplo_leads.csv",
                    mime="text/csv",
                    help="Descarga un archivo CSV de ejemplo de leads",
                    key="leads_ejemplo"
                )
    
    # Tab 4: Matrículas Individuales
    with tabs[3]:
        st.markdown("### Datos Individuales de Matrículas")
        
        # Cargar matrículas actuales si existen
        if matriculas_path.exists():
            df_mats_act = load_dataframe(matriculas_path)
            st.success(f"Matrículas actuales cargadas: {len(df_mats_act)} registros")
            with st.expander("Ver primeros registros"):
                st.dataframe(df_mats_act.head(5))
        else:
            df_mats_act = pd.DataFrame()
            st.info("No hay matrículas actuales cargadas")
        
        # Subir archivo de matrículas
        mats_file = st.file_uploader("Subir archivo de matrículas (CSV/Excel)", type=["csv", "xlsx"], key="matriculas")
        
        if mats_file is not None:
            try:
                if mats_file.name.endswith('.csv'):
                    df_mats_new = pd.read_csv(mats_file)
                else:
                    df_mats_new = pd.read_excel(mats_file)
                
                # Validar estructura
                validation = validate_data_structure(df_mats_new, "matriculas")
                if validation["valid"]:
                    # Verificar duplicados
                    duplicados = check_for_duplicates(df_mats_new)
                    if not duplicados.empty:
                        st.warning(f"Se encontraron {len(duplicados)} IDs duplicados en el archivo.")
                        with st.expander("Ver duplicados"):
                            st.dataframe(duplicados)
                    
                    # Asegurarse que la marca es correcta
                    if "marca" in df_mats_new.columns:
                        df_mats_new["marca"] = brand
                    
                    # Guardar según el modo seleccionado
                    if st.session_state.get("carga_acumulativa", True):
                        # Combinar con matrículas actuales si existen
                        if not df_mats_act.empty:
                            # Eliminar duplicados basados en ID
                            if "ID" in df_mats_act.columns and "ID" in df_mats_new.columns:
                                df_mats_act = df_mats_act[~df_mats_act["ID"].isin(df_mats_new["ID"])]
                            df_mats_act = pd.concat([df_mats_act, df_mats_new], ignore_index=True)
                        else:
                            df_mats_act = df_mats_new
                        
                        # Guardar matrículas actualizadas
                        save_dataframe(df_mats_act, matriculas_path)
                        st.success(f"Matrículas actualizadas (total: {len(df_mats_act)} registros)")
                    else:
                        # Guardar como archivo semanal
                        save_weekly_data(df_mats_new, brand, semana_actual, year_actual, "matriculas")
                        st.success(f"Matrículas guardadas en archivo semanal (Semana {semana_actual}, {year_actual})")
                else:
                    st.error(validation["message"])
                    st.info("Columnas requeridas: ID, fecha_matricula, canal, marca, programa")
            except Exception as e:
                st.error(f"Error al cargar matrículas: {str(e)}")
        
        # Ejemplo de matrículas
        ejemplo_mats = Path("datos/plantillas/ejemplo_matriculas.csv")
        if ejemplo_mats.exists():
            with open(ejemplo_mats, "r") as f:
                st.download_button(
                    label="Ejemplo de estructura de matrículas",
                    data=f,
                    file_name="ejemplo_matriculas.csv",
                    mime="text/csv",
                    help="Descarga un archivo CSV de ejemplo de matrículas",
                    key="mats_ejemplo"
                )
    
    # Preparar DataFrame para uso en reportes
    # Generar DataFrame histórico combinado para compatibilidad con funciones existentes
    df_hist = pd.DataFrame()
    
    # Cargar el histórico si existe
    if hist_path.exists():
        df_hist = load_dataframe(hist_path)
    elif not st.session_state.get("carga_acumulativa", True):
        # Si estamos en modo semanal, combinar todos los archivos semanales
        historico_files = []
        for week in range(1, 53):
            week_file = get_weekly_filename(brand, week, year_actual, "historico")
            week_path = brand_path / HISTORICO_DIR_NAME / week_file
            if week_path.exists():
                historico_files.append(load_dataframe(week_path))
        
        if historico_files:
            df_hist = pd.concat(historico_files, ignore_index=True)
    
    # Si no hay histórico pero hay datos actuales, combinarlos
    if df_hist.empty and 'df_leads_act' in locals() and not df_leads_act.empty and 'df_mats_act' in locals() and not df_mats_act.empty:
        # Número de leads
        if 'fecha_generacion' in df_leads_act.columns:
            leads_por_fecha = df_leads_act.groupby('fecha_generacion').size().reset_index(name='leads')
            leads_por_fecha.rename(columns={'fecha_generacion': 'fecha'}, inplace=True)
            
            # Número de matrículas
            if 'fecha_matricula' in df_mats_act.columns:
                mats_por_fecha = df_mats_act.groupby('fecha_matricula').size().reset_index(name='matriculas')
                mats_por_fecha.rename(columns={'fecha_matricula': 'fecha'}, inplace=True)
                
                # Combinar
                df_hist = pd.merge(leads_por_fecha, mats_por_fecha, on='fecha', how='outer').fillna(0)
                
                # Agregar marca y canal por defecto
                df_hist['marca'] = brand
                df_hist['canal'] = 'Agregado'
                
                # Agregar inversión si existe
                if 'inversion_path' in locals() and inversion_path.exists():
                    df_inv_act = load_dataframe(inversion_path)
                    if not df_inv_act.empty and 'fecha' in df_inv_act.columns and 'monto' in df_inv_act.columns:
                        df_hist = pd.merge(df_hist, df_inv_act[['fecha', 'monto']], on='fecha', how='left')
                        df_hist.rename(columns={'monto': 'inversion'}, inplace=True)
                        df_hist['inversion'].fillna(0, inplace=True)
                    else:
                        df_hist['inversion'] = 0
                else:
                    df_hist['inversion'] = 0
    
    # Guardar el DataFrame combinado para compatibilidad
    if not df_hist.empty:
        save_dataframe(df_hist, hist_path)

    return df_plan, df_hist


# ============================================================
# Reportes
# ============================================================


def calcular_metricas_estrategicas(df_plan: pd.DataFrame, df_hist: pd.DataFrame) -> Dict:
    """Calcula métricas clave para el reporte estratégico."""
    if df_plan.empty or df_hist.empty:
        return {}

    total_inv = df_hist["inversion"].sum()
    total_leads = df_hist["leads"].sum()
    total_mats = df_hist["matriculas"].sum()

    cpa = total_inv / total_mats if total_mats else np.nan
    cpl = total_inv / total_leads if total_leads else np.nan

    progreso_leads = total_leads / df_plan["leads_estimados"].iloc[0] if "leads_estimados" in df_plan.columns else np.nan
    progreso_mats = total_mats / df_plan["objetivo_matriculas"].iloc[0] if "objetivo_matriculas" in df_plan.columns else np.nan

    return {
        "CPA": cpa,
        "CPL": cpl,
        "Progreso Leads": progreso_leads,
        "Progreso Matriculas": progreso_mats,
    }


def show_help(key):
    """Muestra un tooltip de ayuda contextual."""
    st.info(TOOLTIPS.get(key, "Información no disponible"))


def reporte_estrategico_ui(df_plan: pd.DataFrame, df_hist: pd.DataFrame):
    st.header("📊 Reporte Estratégico de Marketing")

    if df_plan.empty or df_hist.empty:
        st.warning("Cargue planificación e histórico para generar el reporte.")
        return

    # Obtener configuración de la convocatoria
    config = st.session_state.get("config_convocatoria", {
        "semana_actual": datetime.datetime.now().isocalendar()[1],
        "duracion_convocatoria": 13,
        "semanas_restantes": 7,
        "progreso": 0.5
    })
    
    # Calcular métricas estratégicas
    metricas = calcular_metricas_estrategicas(df_plan, df_hist)

    # Mostrar información de la convocatoria
    st.subheader("Información de la convocatoria")
    col1, col2, col3, col4 = st.columns(4)
    col1.metric("Semana actual", f"{config['semana_actual']} de {config['duracion_convocatoria']}")
    col2.metric("Semanas restantes", f"{config['semanas_restantes']}")
    col3.metric("Progreso", f"{config['progreso']*100:.1f}%")
    
    # Calcular días transcurridos para reportes
    dias_transcurridos = int(config['progreso'] * config['duracion_convocatoria'] * 7)  # Aproximado en días
    tiempo_total = config['duracion_convocatoria'] * 7  # Total en días
    col4.metric("Días", f"{dias_transcurridos} de {tiempo_total}")

    # Métricas principales
    st.subheader("Métricas principales")
    col1, col2, col3, col4 = st.columns(4)
    col1.metric("CPA", f"${metricas['CPA']:.2f}")
    col2.metric("CPL", f"${metricas['CPL']:.2f}")
    col3.metric("Progreso Leads", f"{metricas['Progreso Leads']*100:.1f}%")
    col4.metric("Progreso Matrículas", f"{metricas['Progreso Matriculas']*100:.1f}%")

    # ======================================================
    # 1. Comparativa de plataformas (CPA y Conversión)
    # ======================================================
    st.subheader("Comparativa de plataformas")

    if "canal" in df_hist.columns and "matriculas" in df_hist.columns and "inversion" in df_hist.columns:
        # Calcular CPA por canal
        df_canales = df_hist.groupby("canal").agg({
            "matriculas": "sum",
            "inversion": "sum",
            "leads": "sum"
        }).reset_index()
        
        # Evitar división por cero
        df_canales["cpa"] = df_canales.apply(
            lambda row: row["inversion"] / row["matriculas"] if row["matriculas"] > 0 else 0, 
            axis=1
        )
        
        df_canales["conversion"] = df_canales.apply(
            lambda row: row["matriculas"] / row["leads"] if row["leads"] > 0 else 0, 
            axis=1
        )
        
        cpa_canal = df_canales.set_index("canal")["cpa"].sort_values()
        conv_canal = df_canales.set_index("canal")["conversion"].sort_values(ascending=False)
        
        st.subheader("CPA por canal")
        st.bar_chart(cpa_canal, use_container_width=True)

        st.subheader("Conversión por canal")
        st.bar_chart(conv_canal, use_container_width=True)
    
    elif "canal" in df_hist.columns and "cpa" in df_hist.columns:
        cpa_canal = df_hist.groupby("canal")["cpa"].mean().sort_values()
        st.bar_chart(cpa_canal, use_container_width=True)

    if "canal" in df_hist.columns and "conversion" in df_hist.columns:
        conv_canal = df_hist.groupby("canal")["conversion"].mean().sort_values()
        st.bar_chart(conv_canal, use_container_width=True)

    # ======================================================
    # 2. Curva de avance: Leads, Matrículas, Inversión
    # ======================================================
    st.subheader("Curva de avance")
    if "fecha" in df_hist.columns:
        df_hist["fecha"] = pd.to_datetime(df_hist["fecha"])
        df_curve = df_hist.sort_values("fecha").set_index("fecha")[[c for c in ["leads", "matriculas", "inversion"] if c in df_hist.columns]]
        st.line_chart(df_curve)
    else:
        st.info("No se encontró la columna 'fecha' para generar curvas.")

    # ======================================================
    # 3. Escenarios simulados basados en tiempo restante
    # ======================================================
    st.subheader("Escenarios simulados")

    # Calcular valores base usando tiempo y configuración
    ratio_tiempo = config['progreso']  # Progreso de tiempo
    
    # Leads y matrículas actuales
    leads_actuales = df_hist["leads"].sum() if "leads" in df_hist.columns else 0
    mats_actuales = df_hist["matriculas"].sum() if "matriculas" in df_hist.columns else 0
    
    # Proyecciones finales (proporcionales al tiempo restante)
    base_leads = leads_actuales / ratio_tiempo if ratio_tiempo > 0 else 0
    base_mats = mats_actuales / ratio_tiempo if ratio_tiempo > 0 else 0
    
    # Objetivos
    leads_obj = df_plan["leads_estimados"].iloc[0] if "leads_estimados" in df_plan.columns else 0
    mats_obj = df_plan["objetivo_matriculas"].iloc[0] if "objetivo_matriculas" in df_plan.columns else 0
    
    # Ratio de conversión actual
    ratio_conversion = mats_actuales / leads_actuales if leads_actuales > 0 else 0
    
    # Proyecciones de escenarios
    escenarios = pd.DataFrame({
        "Escenario": ["Actual", "Mejora 5% conversión", "Aumento +20% inversión"],
        "Leads estimados": [base_leads, base_leads, base_leads*1.2],
        "Matrículas estimadas": [base_mats, base_mats*1.05, base_mats*1.2],
        "% Objetivo": [
            f"{(base_mats/mats_obj)*100:.1f}%" if mats_obj > 0 else "N/A", 
            f"{(base_mats*1.05/mats_obj)*100:.1f}%" if mats_obj > 0 else "N/A", 
            f"{(base_mats*1.2/mats_obj)*100:.1f}%" if mats_obj > 0 else "N/A"
        ],
        "Descripción": [
            "Mantener ritmo actual", 
            "Incremento de conversión en 5% con misma inversión", 
            "Aumentar inversión total 20% manteniendo conversión"
        ]
    })
    st.table(escenarios)

    # ======================================================
    # 4. Predicción ML considerando semanas restantes
    # ======================================================
    cols = st.columns([3, 1])
    cols[0].subheader("Predicción ML de matrículas para próximas semanas")
    if cols[1].button("❓ Ayuda ML", key="help_ml"):
        show_help("prediccion_ml")

    # Usar el número real de semanas restantes para la predicción
    semanas_restantes = config["semanas_restantes"]
    
    brand_path = get_brand_path(st.session_state.current_brand)
    model = train_or_load_model(df_hist, brand_path)
    if model is not None and "leads" in df_hist.columns and "inversion" in df_hist.columns:
        # Crear df_future con promedio de leads e inversión de últimas 4 semanas o todos los disponibles
        avg_leads = df_hist["leads"].tail(min(4, len(df_hist))).mean()
        avg_inv = df_hist["inversion"].tail(min(4, len(df_hist))).mean()
        
        # Usar el número de semanas restantes para la predicción (máx 6 semanas)
        num_semanas_pred = min(6, semanas_restantes)
        
        if num_semanas_pred > 0:
            semanas = pd.DataFrame({
                "leads": [avg_leads] * num_semanas_pred,
                "inversion": [avg_inv] * num_semanas_pred,
                "semana": list(range(1, num_semanas_pred + 1))
            })
            preds, (low, up) = predict_matriculas_interval(model, semanas)
            df_pred = pd.DataFrame({
                "Semana +": list(range(1, num_semanas_pred + 1)),
                "Predicción": preds.astype(int),
                "IC Inferior": low.astype(int),
                "IC Superior": up.astype(int)
            })
            st.dataframe(df_pred)
        else:
            st.info("No quedan semanas restantes para predecir.")
    else:
        st.info("Modelo predictivo no disponible: cargue columnas leads, inversion y matriculas suficientes.")

    # ======================================================
    # 5. Módulo de Atribución Multicanal (NUEVO)
    # ======================================================
    st.subheader("Análisis de Atribución Multicanal")
    cols = st.columns([3, 1])
    cols[0].write("Compare cómo diferentes modelos asignan el mérito a cada canal")
    if cols[1].button("❓ Ayuda Atribución", key="help_atribucion"):
        show_help("atribucion")
    
    # Cargar leads y matrículas individuales
    brand_path = get_brand_path(st.session_state.current_brand)
    leads_path = brand_path / ACTUAL_DIR_NAME / "leads_actual.csv"
    matriculas_path = brand_path / ACTUAL_DIR_NAME / "matriculas_actual.csv"
    
    df_leads = load_dataframe(leads_path)
    df_matriculas = load_dataframe(matriculas_path)
    
    if not df_leads.empty and not df_matriculas.empty:
        # Selector de modelo de atribución
        modelo_atribucion = st.selectbox(
            "Seleccione modelo de atribución:",
            [
                "Último clic (último contacto recibe 100% del mérito)",
                "Primer clic (primer contacto recibe 100% del mérito)",
                "Lineal (mérito distribuido equitativamente)",
                "Decaimiento temporal (más peso a contactos recientes)",
                "Posicional (40% primer, 40% último, 20% intermedios)",
                "Shapley value (basado en contribución marginal)"
            ],
            index=0
        )
        
        # Mapeo de opciones a modelos
        modelo_map = {
            "Último clic (último contacto recibe 100% del mérito)": "ultimo_clic",
            "Primer clic (primer contacto recibe 100% del mérito)": "primer_clic",
            "Lineal (mérito distribuido equitativamente)": "lineal",
            "Decaimiento temporal (más peso a contactos recientes)": "tiempo",
            "Posicional (40% primer, 40% último, 20% intermedios)": "posicional",
            "Shapley value (basado en contribución marginal)": "shapley"
        }
        
        # Calcular atribución con el modelo seleccionado
        modelo_seleccionado = modelo_map[modelo_atribucion]
        df_atribucion = calcular_atribucion(df_leads, df_matriculas, modelo_seleccionado)
        
        if not df_atribucion.empty:
            # Mostrar resultados
            st.write(f"Atribución de matrículas según modelo: **{modelo_atribucion}**")
            st.dataframe(df_atribucion)
            
            # Visualizar como gráfico de barras
            st.bar_chart(df_atribucion.set_index("canal")["atribucion"])
            
            # Opción para comparar todos los modelos
            if st.checkbox("Comparar todos los modelos de atribución"):
                todos_resultados = comparar_modelos_atribucion(df_leads, df_matriculas)
                
                # Crear tabla comparativa
                tabla_comparativa = pd.DataFrame()
                
                for modelo, df in todos_resultados.items():
                    if not df.empty:
                        # Tomar solo los 3 principales canales para cada modelo
                        top_canales = df.head(3)[["canal", "atribucion"]]
                        
                        # Renombrar columnas
                        top_canales = top_canales.copy()
                        top_canales["modelo"] = modelo
                        top_canales.rename(columns={"atribucion": "valor"}, inplace=True)
                        
                        # Agregar a la tabla comparativa
                        tabla_comparativa = pd.concat([tabla_comparativa, top_canales], ignore_index=True)
                
                if not tabla_comparativa.empty:
                    # Crear vista pivot para comparar modelos
                    pivot = tabla_comparativa.pivot_table(
                        index="canal", 
                        columns="modelo", 
                        values="valor", 
                        aggfunc="sum"
                    ).fillna(0)
                    
                    st.write("Comparativa de atribución por modelo (top canales)")
                    st.dataframe(pivot)
        else:
            st.warning("No se pudo generar el análisis de atribución. Verifique los datos.")
    else:
        st.info("Cargue datos individuales de leads y matrículas para utilizar el análisis de atribución multicanal.")

    # ======================================================
    # 6. Alertas automáticas basadas en tiempo restante
    # ======================================================
    st.subheader("Alertas destacadas")
    
    # Verificar CPA alto
    if metricas["CPA"] > 50:
        st.error("CPA alto: revise campañas menos eficientes.")
    
    # Verificar progreso considerando tiempo restante
    progreso_leads = metricas["Progreso Leads"]
    tiempo_usado = config["progreso"]
    
    # Verificar si el progreso de leads está por debajo del progreso de tiempo
    if progreso_leads < tiempo_usado * 0.8:
        st.warning(f"Ritmo bajo de generación de leads ({progreso_leads*100:.1f}% completado vs {tiempo_usado*100:.1f}% de tiempo usado).")
    
    # Verificar canales ineficientes
    if "cpa_canal" in locals():
        inefficient = cpa_canal[cpa_canal > metricas["CPA"] * 1.2]
        for canal, val in inefficient.items():
            st.warning(f"Canal ineficiente: {canal} (CPA ${val:.2f})")
    
    # Alerta por tiempo restante limitado
    if semanas_restantes <= 4 and metricas["Progreso Matriculas"] < 0.6:
        st.error(f"¡URGENTE! Quedan solo {semanas_restantes} semanas y se ha alcanzado solo el {metricas['Progreso Matriculas']*100:.1f}% del objetivo de matrículas.")

    # ======================================================
    # 7. Exportación a Excel mejorada
    # ======================================================
    st.subheader("Exportar informe")
    if st.button("Descargar Excel Avanzado", key="excel_estrategico_avanzado"):
        from io import BytesIO
        buffer = BytesIO()
        
        # Mensaje de procesamiento
        with st.spinner("Generando Excel optimizado..."):
            with pd.ExcelWriter(buffer, engine="xlsxwriter") as writer:
                workbook = writer.book
                
                # Crear formatos atractivos
                header_format = workbook.add_format({
                    'bold': True,
                    'text_wrap': True,
                    'valign': 'top',
                    'fg_color': '#1E3A8A',
                    'font_color': 'white',
                    'border': 1,
                    'align': 'center'
                })
                
                number_format = workbook.add_format({
                    'num_format': '#,##0',
                    'border': 1,
                    'align': 'center'
                })
                
                currency_format = workbook.add_format({
                    'num_format': '$#,##0.00',
                    'border': 1,
                    'align': 'center'
                })
                
                percent_format = workbook.add_format({
                    'num_format': '0.00%',
                    'border': 1,
                    'align': 'center'
                })
                
                title_format = workbook.add_format({
                    'bold': True, 
                    'font_size': 14, 
                    'align': 'center',
                    'valign': 'vcenter',
                    'font_color': '#1E3A8A'
                })
                
                subtitle_format = workbook.add_format({
                    'bold': True, 
                    'font_size': 12, 
                    'align': 'center',
                    'italic': True,
                    'font_color': '#4B5563'
                })
                
                border_format = workbook.add_format({
                    'border': 1
                })
                
                # Hoja 1: Resumen Ejecutivo
                # --------------------------
                resumen_data = {
                    "Métrica": ["Fecha Reporte", "Marca", "Semana Actual", "Progreso Tiempo", 
                               "Leads Obtenidos", "Meta Leads", "Matrículas Actuales", "Meta Matrículas",
                               "CPA Actual", "CPL Actual", "Conversión Actual"],
                    "Valor": [datetime.datetime.now().strftime("%d/%m/%Y"), 
                             st.session_state.current_brand,
                             f"{config['semana_actual']} de {config['duracion_convocatoria']}",
                             f"{config['progreso']*100:.1f}%",
                             df_hist["leads"].sum(),
                             leads_obj,
                             df_hist["matriculas"].sum(),
                             mats_obj,
                             f"${metricas['CPA']:.2f}",
                             f"${metricas['CPL']:.2f}",
                             f"{mats_actuales/leads_actuales*100:.2f}%" if leads_actuales > 0 else "0%"
                            ]
                }
                
                df_resumen = pd.DataFrame(resumen_data)
                
                # Primer hoja - Resumen Ejecutivo
                df_resumen.to_excel(writer, sheet_name='Resumen Ejecutivo', index=False, startrow=2)
                worksheet = writer.sheets['Resumen Ejecutivo']
                
                # Agregar título
                worksheet.merge_range('A1:B1', f"REPORTE ESTRATÉGICO - {st.session_state.current_brand}", title_format)
                worksheet.merge_range('A2:B2', f"Generado: {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}", subtitle_format)
                
                # Aplicar formato a columnas
                worksheet.set_column('A:A', 20)
                worksheet.set_column('B:B', 25)
                
                # Aplicar estilos a cabeceras
                for col_num, value in enumerate(df_resumen.columns.values):
                    worksheet.write(2, col_num, value, header_format)
                
                # Hoja 2: Datos Históricos con formato
                # -----------------------------------
                if not df_hist.empty:
                    df_hist.to_excel(writer, sheet_name="Datos Históricos", index=False, startrow=1)
                    worksheet = writer.sheets['Datos Históricos']
                    
                    # Aplicar título
                    worksheet.merge_range('A1:' + chr(65 + len(df_hist.columns) - 1) + '1', 
                                         "DATOS HISTÓRICOS COMPLETOS", title_format)
                    
                    # Aplicar formato a cabeceras
                    for col_num, value in enumerate(df_hist.columns.values):
                        worksheet.write(1, col_num, value, header_format)
                    
                    # Aplicar formato a columnas numéricas
                    for i, col in enumerate(df_hist.columns):
                        # Ajustar ancho automáticamente
                        max_len = max(df_hist[col].astype(str).str.len().max(), len(col))
                        worksheet.set_column(i, i, max_len + 3)
                        
                        # Aplicar formatos según tipo de columna
                        if col in ['inversion', 'CPA', 'CPL']:
                            worksheet.set_column(i, i, max_len + 3, currency_format)
                        elif col in ['leads', 'matriculas']:
                            worksheet.set_column(i, i, max_len + 3, number_format)
                        elif 'porcentaje' in col.lower() or 'tasa' in col.lower() or 'conversion' in col.lower():
                            worksheet.set_column(i, i, max_len + 3, percent_format)
                    
                    # Agregar filtros automáticos
                    worksheet.autofilter(1, 0, len(df_hist) + 1, len(df_hist.columns) - 1)
                
                # Hoja 3: Métricas por Canal
                # -------------------------
                if "df_canales" in locals() and not df_canales.empty:
                    df_canales.to_excel(writer, sheet_name="Métricas por Canal", index=False, startrow=1)
                    worksheet = writer.sheets['Métricas por Canal']
                    
                    # Título
                    worksheet.merge_range('A1:' + chr(65 + len(df_canales.columns) - 1) + '1', 
                                         "ANÁLISIS POR CANAL DE MARKETING", title_format)
                    
                    # Cabeceras
                    for col_num, value in enumerate(df_canales.columns.values):
                        worksheet.write(1, col_num, value, header_format)
                    
                    # Formato a columnas
                    for i, col in enumerate(df_canales.columns):
                        max_len = max(df_canales[col].astype(str).str.len().max(), len(col))
                        worksheet.set_column(i, i, max_len + 3)
                        
                        if col in ['inversion', 'cpa']:
                            worksheet.set_column(i, i, max_len + 3, currency_format)
                        elif col in ['leads', 'matriculas']:
                            worksheet.set_column(i, i, max_len + 3, number_format)
                        elif col == 'conversion':
                            worksheet.set_column(i, i, max_len + 3, percent_format)
                
                # Hoja 4: Escenarios Proyectados
                # -----------------------------
                escenarios.to_excel(writer, sheet_name="Escenarios", index=False, startrow=1)
                worksheet = writer.sheets['Escenarios']
                
                # Título
                worksheet.merge_range('A1:' + chr(65 + len(escenarios.columns) - 1) + '1', 
                                     "PROYECCIÓN DE ESCENARIOS", title_format)
                
                # Formato cabeceras
                for col_num, value in enumerate(escenarios.columns.values):
                    worksheet.write(1, col_num, value, header_format)
                
                # Ajustar columnas
                for i, col in enumerate(escenarios.columns):
                    max_len = max(escenarios[col].astype(str).str.len().max(), len(col))
                    worksheet.set_column(i, i, max_len + 3)
                
                # Hoja 5: Atribución (si existe)
                # ----------------------------
                if "df_atribucion" in locals() and not df_atribucion.empty:
                    df_atribucion.to_excel(writer, sheet_name="Atribución", index=False, startrow=1)
                    worksheet = writer.sheets['Atribución']
                    
                    # Título
                    worksheet.merge_range('A1:' + chr(65 + len(df_atribucion.columns) - 1) + '1', 
                                         f"ATRIBUCIÓN MULTICANAL - MODELO {modelo_seleccionado.upper()}", title_format)
                    
                    # Cabeceras
                    for col_num, value in enumerate(df_atribucion.columns.values):
                        worksheet.write(1, col_num, value, header_format)
                    
                    # Formato columnas
                    for i, col in enumerate(df_atribucion.columns):
                        max_len = max(df_atribucion[col].astype(str).str.len().max(), len(col))
                        worksheet.set_column(i, i, max_len + 3)
                        
                        if col == 'porcentaje':
                            worksheet.set_column(i, i, max_len + 3, percent_format)
                    
                    # Añadimos un gráfico de pastel para atribución
                    chart_atrib = workbook.add_chart({'type': 'pie'})
                    
                    # Configurar los datos del gráfico (ajustando filas)
                    chart_atrib.add_series({
                        'name': 'Atribución por Canal',
                        'categories': ['Atribución', 2, 0, 1 + len(df_atribucion), 0],
                        'values': ['Atribución', 2, 1, 1 + len(df_atribucion), 1],
                        'data_labels': {'percentage': True, 'position': 'outside_end'},
                    })
                    
                    # Configurar título y posición
                    chart_atrib.set_title({'name': 'Distribución de Atribución por Canal'})
                    chart_atrib.set_style(10)
                    worksheet.insert_chart('F2', chart_atrib, {'x_offset': 25, 'y_offset': 10, 'x_scale': 1.5, 'y_scale': 1.5})
                
                # Hoja 6: Predicciones ML (si existe)
                # ---------------------------------
                if "df_pred" in locals() and not df_pred.empty:
                    df_pred.to_excel(writer, sheet_name="Predicciones ML", index=False, startrow=1)
                    worksheet = writer.sheets['Predicciones ML']
                    
                    # Título
                    worksheet.merge_range('A1:' + chr(65 + len(df_pred.columns) - 1) + '1', 
                                         "PREDICCIONES MACHINE LEARNING", title_format)
                    
                    # Cabeceras
                    for col_num, value in enumerate(df_pred.columns.values):
                        worksheet.write(1, col_num, value, header_format)
                    
                    # Formato columnas
                    for i, col in enumerate(df_pred.columns):
                        worksheet.set_column(i, i, 15, number_format)
                    
                    # Crear gráfico de líneas para predicciones
                    chart = workbook.add_chart({'type': 'line'})
                    
                    # Predicción central
                    chart.add_series({
                        'name': 'Predicción',
                        'categories': ['Predicciones ML', 2, 0, 1 + len(df_pred), 0],
                        'values': ['Predicciones ML', 2, 1, 1 + len(df_pred), 1],
                        'line': {'color': 'blue', 'width': 2.5},
                        'marker': {'type': 'circle', 'size': 5, 'fill': {'color': 'blue'}}
                    })
                    
                    # Límite inferior
                    chart.add_series({
                        'name': 'Límite Inferior',
                        'categories': ['Predicciones ML', 2, 0, 1 + len(df_pred), 0],
                        'values': ['Predicciones ML', 2, 2, 1 + len(df_pred), 2],
                        'line': {'color': 'orange', 'width': 1, 'dash_type': 'dash'},
                    })
                    
                    # Límite superior
                    chart.add_series({
                        'name': 'Límite Superior',
                        'categories': ['Predicciones ML', 2, 0, 1 + len(df_pred), 0],
                        'values': ['Predicciones ML', 2, 3, 1 + len(df_pred), 3],
                        'line': {'color': 'green', 'width': 1, 'dash_type': 'dash'},
                    })
                    
                    # Configurar gráfico
                    chart.set_title({'name': 'Proyección de Matrículas'})
                    chart.set_x_axis({'name': 'Semana +'})
                    chart.set_y_axis({'name': 'Matrículas Proyectadas'})
                    chart.set_style(42)
                    worksheet.insert_chart('F2', chart, {'x_offset': 25, 'y_offset': 10, 'x_scale': 1.5, 'y_scale': 1.5})
            
            # Botón de descarga
            fecha_actual = datetime.datetime.now().strftime("%Y%m%d_%H%M")
            buffer.seek(0)
            st.download_button(
                label="⬇️ Descargar Excel Completo",
                data=buffer,
                file_name=f"reporte_estrategico_{st.session_state.current_brand}_{fecha_actual}.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                key="excel_estrategico_btn"
            )
            st.success("✅ Reporte Excel generado correctamente con formatos optimizados y gráficos interactivos.")


def reporte_comercial_ui(df_plan: pd.DataFrame, df_hist: pd.DataFrame):
    st.header("📊 Reporte Comercial (Status Semanal)")

    if df_plan.empty or df_hist.empty:
        st.warning("Cargue planificación e histórico para generar el reporte.")
        return

    # Obtener configuración de la convocatoria
    config = st.session_state.get("config_convocatoria", {
        "semana_actual": datetime.datetime.now().isocalendar()[1],
        "duracion_convocatoria": 13,
        "semanas_restantes": 7,
        "progreso": 0.5
    })
    
    semana_actual = config["semana_actual"]
    duracion_total = config["duracion_convocatoria"]
    semanas_restantes = config["semanas_restantes"]
    progreso = config["progreso"]

    # Fechas
    if "fecha" in df_hist.columns:
        df_hist["fecha"] = pd.to_datetime(df_hist["fecha"])
        fecha_min, fecha_max = df_hist["fecha"].min(), df_hist["fecha"].max()
        dias_transcurridos = (fecha_max - fecha_min).days + 1
    else:
        dias_transcurridos = int(progreso * duracion_total * 7)  # Aproximado en días

    # Objetivos
    tiempo_total = duracion_total * 7  # Días totales de la convocatoria
    leads_obj = df_plan["leads_estimados"].iloc[0] if "leads_estimados" in df_plan.columns else 0
    mats_obj = df_plan["objetivo_matriculas"].iloc[0] if "objetivo_matriculas" in df_plan.columns else 0

    # Valores actuales
    leads_actuales = df_hist["leads"].sum() if "leads" in df_hist.columns else 0
    mats_actuales = df_hist["matriculas"].sum() if "matriculas" in df_hist.columns else 0

    # Información de la convocatoria
    st.subheader("Información de la convocatoria")
    col1, col2, col3 = st.columns(3)
    col1.metric("Semana actual", f"{semana_actual} de {duracion_total}")
    col2.metric("Semanas restantes", f"{semanas_restantes}")
    col3.metric("Progreso", f"{progreso*100:.1f}%")

    # Barras de progreso mejoradas
    st.subheader("Progreso de la convocatoria")
    
    # Función para crear barras de progreso con colores
    def barra_color(texto, valor, total):
        pct = int(100 * valor / total) if total else 0
        # Determinar color según porcentaje
        if pct < 60:
            color = "🔴"
        elif pct < 90:
            color = "🟠"
        else:
            color = "🟢"
        
        bar = "▓" * (pct // 10) + "░" * (10 - pct // 10)
        st.write(f"{texto:<25} {bar} {color} {pct}%")

    barra_color("Tiempo transcurrido", dias_transcurridos, tiempo_total)
    barra_color("Leads entregados", leads_actuales, leads_obj)
    barra_color("Matrículas confirmadas", mats_actuales, mats_obj)

    # Proyección simple lineal basada en el tiempo transcurrido
    ratio_tiempo = progreso  # Progreso de tiempo
    ratio_conversión = mats_actuales / leads_actuales if leads_actuales > 0 else 0
    
    # Proyección de leads
    proy_leads_final = leads_actuales / ratio_tiempo if ratio_tiempo > 0 else 0
    
    # Proyección de matrículas basada en proyección de leads y ratio de conversión
    proy_mats_por_leads = proy_leads_final * ratio_conversión
    
    # Proyección de matrículas basada en ritmo actual
    proy_mats_por_tiempo = mats_actuales / ratio_tiempo if ratio_tiempo > 0 else 0
    
    # Proyección final (promedio de ambos métodos)
    proy_final = (proy_mats_por_leads + proy_mats_por_tiempo) / 2
    
    # Intervalo de confianza usando tiempo restante
    margen_error = proy_final * (semanas_restantes / duracion_total) * 0.2  # Margen de error proporcional al tiempo restante
    ic_lower = max(0, proy_final - margen_error)
    ic_upper = proy_final + margen_error
    
    # Determinar estado basado en la proyección vs objetivo
    if proy_final >= mats_obj:
        estado = "normal"
        estado_texto = "✅ Normal (en camino al objetivo)"
    elif proy_final >= mats_obj * 0.8:
        estado = "leve_desvio"
        estado_texto = "⚠️ Leve desvío (proyección cercana al objetivo)"
    else:
        estado = "riesgo_alto"
        estado_texto = "🔴 Riesgo alto (proyección muy por debajo del objetivo)"
    
    col1, col2 = st.columns([3, 1])
    col1.metric("Proyección matrículas finales", f"{int(proy_final)}", delta=f"±{int(margen_error)} (IC 95%)")
    col2.metric("Estado", estado_texto)
    
    # Visualización de tendencia
    st.subheader("Tendencia de Cumplimiento")
    # Calcular tendencia basada en datos semanales
    if "fecha" in df_hist.columns and len(df_hist) > 1:
        df_hist["fecha"] = pd.to_datetime(df_hist["fecha"])
        df_hist = df_hist.sort_values("fecha")
        
        # Agrupar por semana
        df_hist["dia_semana"] = df_hist["fecha"].dt.day_name()
        df_hist["mes"] = df_hist["fecha"].dt.month_name()
        
        # Análisis por día de semana
        leads_dia = df_hist.groupby("dia_semana")["leads"].mean().sort_values(ascending=False)
        st.subheader("Días con mayor captación de leads")
        st.bar_chart(leads_dia)
        
        # Mejor día
        mejor_dia = leads_dia.idxmax()
        st.info(f"El día con mayor captación promedio de leads es {mejor_dia}")
        
        # Análisis de tendencia actual vs histórico
        st.subheader("Análisis de tendencia en convocatoria actual")
        
        # Obtener fechas de inicio y fin basadas en la configuración
        duracion_dias = config['duracion_convocatoria'] * 7
        dias_transcurridos = int(config['progreso'] * duracion_dias)
        
        if dias_transcurridos > 0:
            fechas_actuales = pd.date_range(end=pd.Timestamp.now(), periods=min(dias_transcurridos, len(df_hist)))
            
            # Filtrar datos solo para la convocatoria actual
            df_hist_actual = df_hist[df_hist["fecha"] >= fechas_actuales[0]]
            if not df_hist_actual.empty:
                st.line_chart(df_hist_actual.set_index("fecha")["leads"])
                
                # Comparar con tendencia histórica
                media_leads_dia_actual = df_hist_actual["leads"].mean()
                media_leads_dia_historica = df_hist["leads"].mean()
                
                tendencia = media_leads_dia_actual / media_leads_dia_historica
                if tendencia > 1.1:
                    st.success(f"Tendencia positiva: +{(tendencia-1)*100:.1f}% leads/día vs histórico")
                elif tendencia < 0.9:
                    st.warning(f"Tendencia negativa: {(tendencia-1)*100:.1f}% leads/día vs histórico")
                else:
                    st.info(f"Tendencia estable: {tendencia*100:.1f}% vs histórico")
    else:
        st.info("Se requiere columna 'fecha' para realizar análisis temporal.")

    # Exportar
    st.subheader("Exportar informe")
    if st.button("Descargar Excel Avanzado", key="excel_comercial_avanzado"):
        from io import BytesIO
        buffer = BytesIO()
        
        # Mensaje de procesamiento
        with st.spinner("Generando Excel comercial optimizado..."):
            with pd.ExcelWriter(buffer, engine="xlsxwriter") as writer:
                workbook = writer.book
                
                # Crear formatos atractivos para Excel
                header_format = workbook.add_format({
                    'bold': True,
                    'text_wrap': True,
                    'valign': 'top',
                    'fg_color': '#1E3A8A',
                    'font_color': 'white',
                    'border': 1,
                    'align': 'center'
                })
                
                number_format = workbook.add_format({
                    'num_format': '#,##0',
                    'border': 1,
                    'align': 'center'
                })
                
                currency_format = workbook.add_format({
                    'num_format': '$#,##0.00',
                    'border': 1,
                    'align': 'center'
                })
                
                percent_format = workbook.add_format({
                    'num_format': '0.00%',
                    'border': 1,
                    'align': 'center'
                })
                
                title_format = workbook.add_format({
                    'bold': True, 
                    'font_size': 14, 
                    'align': 'center',
                    'valign': 'vcenter',
                    'font_color': '#1E3A8A'
                })
                
                subtitle_format = workbook.add_format({
                    'bold': True, 
                    'font_size': 12, 
                    'align': 'center',
                    'italic': True,
                    'font_color': '#4B5563'
                })
                
                border_format = workbook.add_format({
                    'border': 1
                })
                
                # Formato para barras de progreso
                progress_red = workbook.add_format({
                    'bg_color': '#FFAAAA',
                    'pattern': 1,
                    'border': 1
                })
                
                progress_yellow = workbook.add_format({
                    'bg_color': '#FFFFAA',
                    'pattern': 1,
                    'border': 1
                })
                
                progress_green = workbook.add_format({
                    'bg_color': '#AAFFAA',
                    'pattern': 1,
                    'border': 1
                })
                
                # Hoja 1: Resumen General
                # ----------------------
                resumen_data = {
                    "Métrica": ["Fecha Reporte", "Marca", "Semana Actual", "Semanas Restantes", 
                               "Progreso Tiempo", "Leads Generados", "Meta Leads", 
                               "Matrículas Actuales", "Meta Matrículas",
                               "Estado Campaña", "Proyección Matrículas"],
                    "Valor": [datetime.datetime.now().strftime("%d/%m/%Y"), 
                             st.session_state.current_brand,
                             f"{semana_actual} de {duracion_total}",
                             f"{semanas_restantes}",
                             f"{progreso*100:.1f}%",
                             leads_actuales,
                             leads_obj,
                             mats_actuales,
                             mats_obj,
                             estado_texto,
                             f"{int(proy_final)} (±{int(margen_error)})"
                            ]
                }
                
                df_resumen = pd.DataFrame(resumen_data)
                
                # Escribir a Excel
                df_resumen.to_excel(writer, sheet_name='Resumen', index=False, startrow=2)
                worksheet = writer.sheets['Resumen']
                
                # Agregar título
                worksheet.merge_range('A1:B1', f"REPORTE COMERCIAL SEMANAL - {st.session_state.current_brand}", title_format)
                worksheet.merge_range('A2:B2', f"Generado: {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}", subtitle_format)
                
                # Aplicar formato a columnas
                worksheet.set_column('A:A', 20)
                worksheet.set_column('B:B', 25)
                
                # Aplicar estilos a cabeceras
                for col_num, value in enumerate(df_resumen.columns.values):
                    worksheet.write(2, col_num, value, header_format)
                
                # Hoja 2: Datos Históricos
                # -----------------------
                if not df_hist.empty:
                    df_hist.to_excel(writer, sheet_name="Datos Históricos", index=False, startrow=1)
                    worksheet = writer.sheets['Datos Históricos']
                    
                    # Aplicar título
                    worksheet.merge_range('A1:' + chr(65 + len(df_hist.columns) - 1) + '1', 
                                         "DATOS HISTÓRICOS COMPLETOS", title_format)
                    
                    # Aplicar formato a cabeceras
                    for col_num, value in enumerate(df_hist.columns.values):
                        worksheet.write(1, col_num, value, header_format)
                    
                    # Aplicar formato a columnas numéricas
                    for i, col in enumerate(df_hist.columns):
                        # Ajustar ancho automáticamente
                        max_len = max(df_hist[col].astype(str).str.len().max(), len(col))
                        worksheet.set_column(i, i, max_len + 3)
                        
                        # Aplicar formatos según tipo de columna
                        if col in ['inversion', 'CPA', 'CPL']:
                            worksheet.set_column(i, i, max_len + 3, currency_format)
                        elif col in ['leads', 'matriculas']:
                            worksheet.set_column(i, i, max_len + 3, number_format)
                        elif 'porcentaje' in col.lower() or 'tasa' in col.lower() or 'conversion' in col.lower():
                            worksheet.set_column(i, i, max_len + 3, percent_format)
                
                # Hoja 3: Barras de Progreso
                # ------------------------
                barras_data = {
                    "Elemento": ["Tiempo transcurrido", "Leads entregados", "Matrículas confirmadas"],
                    "Valor Actual": [dias_transcurridos, leads_actuales, mats_actuales],
                    "Valor Objetivo": [tiempo_total, leads_obj, mats_obj],
                    "Porcentaje": [
                        dias_transcurridos / tiempo_total if tiempo_total else 0,
                        leads_actuales / leads_obj if leads_obj else 0,
                        mats_actuales / mats_obj if mats_obj else 0
                    ]
                }
                
                df_barras = pd.DataFrame(barras_data)
                df_barras.to_excel(writer, sheet_name="Progreso", index=False, startrow=1)
                worksheet = writer.sheets['Progreso']
                
                # Aplicar título
                worksheet.merge_range('A1:D1', "BARRAS DE PROGRESO", title_format)
                
                # Aplicar formato a cabeceras
                for col_num, value in enumerate(df_barras.columns.values):
                    worksheet.write(1, col_num, value, header_format)
                
                # Ajustar ancho de columnas
                worksheet.set_column('A:A', 20)
                worksheet.set_column('B:B', 15)
                worksheet.set_column('C:C', 15)
                worksheet.set_column('D:D', 15)
                
                # Aplicar formato de porcentaje a la columna de porcentaje
                for i in range(len(df_barras)):
                    worksheet.write(i+2, 3, df_barras["Porcentaje"].iloc[i], percent_format)
                
                # Crear barras de progreso visuales (fila 10 en adelante)
                worksheet.write(9, 0, "Barras de Progreso Visual", subtitle_format)
                
                for i, row in enumerate(df_barras.itertuples()):
                    elem_name = row.Elemento
                    pct = row.Porcentaje
                    
                    # Determinar color basado en porcentaje
                    formato = progress_green if pct >= 0.9 else (progress_yellow if pct >= 0.6 else progress_red)
                    
                    # Escribir nombre de elemento y crear barra
                    worksheet.write(i+10, 0, elem_name)
                    cells_filled = int(pct * 10)  # Barra de 10 celdas
                    
                    # Dibujar progreso como celdas coloreadas
                    for j in range(cells_filled):
                        worksheet.write(i+10, j+1, "", formato)
                    
                    # Escribir porcentaje al final de la barra
                    worksheet.write(i+10, 12, f"{pct*100:.1f}%")
                
                # Hoja 4: Tendencia por Día
                # -----------------------
                if "dia_semana" in locals() and "leads_dia" in locals() and not leads_dia.empty:
                    df_dias = pd.DataFrame({"Día": leads_dia.index, "Promedio Leads": leads_dia.values})
                    df_dias.to_excel(writer, sheet_name="Tendencia por Día", index=False, startrow=1)
                    worksheet = writer.sheets['Tendencia por Día']
                    
                    # Aplicar título
                    worksheet.merge_range('A1:B1', "CAPTACIÓN POR DÍA DE LA SEMANA", title_format)
                    
                    # Aplicar formato a cabeceras
                    for col_num, value in enumerate(df_dias.columns.values):
                        worksheet.write(1, col_num, value, header_format)
                    
                    # Ajustar ancho de columnas
                    worksheet.set_column('A:A', 15)
                    worksheet.set_column('B:B', 15)
                    
                    # Crear gráfico de columnas
                    chart = workbook.add_chart({'type': 'column'})
                    chart.add_series({
                        'name': 'Leads por día',
                        'categories': ['Tendencia por Día', 2, 0, 1 + len(df_dias), 0],
                        'values': ['Tendencia por Día', 2, 1, 1 + len(df_dias), 1],
                        'data_labels': {'value': True}
                    })
                    
                    chart.set_title({'name': 'Captación por Día de la Semana'})
                    chart.set_x_axis({'name': 'Día'})
                    chart.set_y_axis({'name': 'Promedio Leads'})
                    chart.set_style(11)  # Estilo más moderno
                    
                    worksheet.insert_chart('D2', chart, {'x_offset': 25, 'y_offset': 10, 'x_scale': 1.5, 'y_scale': 1.2})
                
                # Hoja 5: Proyecciones
                # -----------------
                proyecciones_data = {
                    "Métrica": ["Leads Actuales", "Leads Proyección Final", 
                               "Matrículas Actuales", "Matrículas Proyección por Leads", 
                               "Matrículas Proyección por Tiempo", "Matrículas Proyección Final",
                               "Límite Inferior", "Límite Superior", "Meta Objetivo"],
                    "Valor": [leads_actuales, proy_leads_final, 
                             mats_actuales, proy_mats_por_leads,
                             proy_mats_por_tiempo, proy_final, 
                             ic_lower, ic_upper, mats_obj]
                }
                
                df_proyecciones = pd.DataFrame(proyecciones_data)
                df_proyecciones.to_excel(writer, sheet_name="Proyecciones", index=False, startrow=1)
                worksheet = writer.sheets['Proyecciones']
                
                # Aplicar título
                worksheet.merge_range('A1:B1', "PROYECCIONES FINALES", title_format)
                
                # Aplicar formato a cabeceras
                for col_num, value in enumerate(df_proyecciones.columns.values):
                    worksheet.write(1, col_num, value, header_format)
                
                # Ajustar ancho de columnas y aplicar formatos
                worksheet.set_column('A:A', 25)
                worksheet.set_column('B:B', 15, number_format)
                
                # Crear gráfico para proyecciones
                chart = workbook.add_chart({'type': 'column'})
                
                # Series para el gráfico de proyecciones
                chart.add_series({
                    'name': 'Actual vs Proyectado',
                    'categories': ['Proyecciones', 3, 0, 4, 0],  # Solo matrículas actuales y finales
                    'values': ['Proyecciones', 3, 1, 4, 1],
                    'data_labels': {'value': True}
                })
                
                # Añadir línea para el objetivo
                chart.add_series({
                    'name': 'Objetivo',
                    'categories': ['Proyecciones', 8, 0, 8, 0],  # Solo la fila de objetivo
                    'values': ['Proyecciones', 8, 1, 8, 1],
                    'type': 'line',
                    'line': {'color': 'red', 'width': 2, 'dash_type': 'dash'},
                    'marker': {'type': 'square', 'size': 8}
                })
                
                chart.set_title({'name': 'Matrículas: Actual vs Proyección vs Objetivo'})
                chart.set_style(42)
                worksheet.insert_chart('D2', chart, {'x_offset': 25, 'y_offset': 10, 'x_scale': 1.5, 'y_scale': 1.2})
                
            # Botón de descarga
            fecha_actual = datetime.datetime.now().strftime("%Y%m%d_%H%M")
            buffer.seek(0)
            st.download_button(
                label="⬇️ Descargar Excel Completo",
                data=buffer,
                file_name=f"reporte_comercial_{st.session_state.current_brand}_{fecha_actual}.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                key="excel_comercial_btn"
            )
            st.success("✅ Reporte comercial Excel generado correctamente con formatos optimizados y gráficos interactivos.")


# ============================================================
# Main App
# ============================================================


def main():
    st.title("📘 Motor de Decisión Educativo y Predictivo")
    
    # Aplicar estilos CSS personalizados
    st.markdown(CUSTOM_CSS, unsafe_allow_html=True)

    # Agregar ayuda general del sistema
    if st.sidebar.button("📋 Manual de usuario"):
        st.sidebar.info("""
        **Motor de Decisión Educativo v{}**

        Este sistema permite analizar datos de marketing educativo por marcas,
        generar informes estratégicos y comerciales, y exportar datos
        para toma de decisiones.

        **Uso básico:**
        1. Seleccione una marca o cree una nueva
        2. Cargue datos de planificación e histórico
        3. Explore los reportes disponibles
        4. Utilice los botones de ayuda (❓) para información contextual

        **Requerimientos:**
        - fpdf (para exportación PDF)
        - python-pptx (para exportación PowerPoint)
        - scikit-learn, joblib (para predicciones ML)
        """.format(VERSION))

    brand = sidebar_brand_selector()
    st.sidebar.write(f"Marca seleccionada: **{brand}**")

    menu = st.sidebar.radio(
        "Menú", ["Carga de Datos", "Reporte Estratégico", "Reporte Comercial", "Reporte Exploratorio"], index=0
    )

    df_plan, df_hist = load_data_ui(brand)

    if menu == "Reporte Estratégico":
        reporte_estrategico_ui(df_plan, df_hist)
    elif menu == "Reporte Comercial":
        reporte_comercial_ui(df_plan, df_hist)
    elif menu == "Reporte Exploratorio":
        reporte_exploratorio_ui(df_plan, df_hist)


if __name__ == "__main__":
    main() 